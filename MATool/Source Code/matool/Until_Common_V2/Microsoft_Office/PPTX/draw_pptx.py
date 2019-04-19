from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE,XL_LEGEND_POSITION,XL_LABEL_POSITION
from pprint import pprint

class DrawPPTX:
	__slide = ''

	def __init__(self,slide):
		self.__slide = slide

	def create_two_lines_text(self,text_1,text_2,left,top,back_ground_color=RGBColor(0, 0, 0),
														height=1,width=3,font_size_text_1=16,font_size_text_2=16,
														text_color=RGBColor(255, 255, 255),font_bold = False,
														font_name = "Arial",alignment=False):
		# 1)Vị trí của đoạn text
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)

		txBox = self.__slide.shapes.add_textbox(left, top, width, height)
		# 1.1)Chỉnh màu background
		txBox.fill.solid()
		txBox.fill.fore_color.	rgb = back_ground_color

		# 2)Tạo các dòng text
		tf = txBox.text_frame
		tf.word_wrap = True
		# 2.1) Dòng 1 trong text_box
		p = tf.paragraphs[0]
		p.text = text_1
		p.font.size = Pt(font_size_text_1)
		p.font.color.rgb = text_color
		p.font.bold = font_bold
		p.font.name = font_name
		#=== Căn giữa ===
		if alignment == True:
			p.alignment = PP_ALIGN.CENTER

		# 2.2) Dòng 2 trong text_box
		p = tf.add_paragraph()
		p.text = text_2
		p.font.size = Pt(font_size_text_2)
		p.font.color.rgb = text_color
		p.font.bold = font_bold
		p.font.name = font_name
		#=== Căn giữa ===
		if alignment == True:
			p.alignment = PP_ALIGN.CENTER
		
	
	def create_one_line_text(self,text_1,left,top,back_ground_color=RGBColor(0, 0, 0),
													height=1,width=3,font_size=16,text_color=RGBColor(255, 255, 255),
													font_bold = False,font_name = "Arial",alignment=False):
		#1)TEXT BOX COLOR
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)

		txBox = self.__slide.shapes.add_textbox(left, top, width, height)
		#1.1)FILL BACK GROUND COLOR
		txBox.fill.solid()
		txBox.fill.fore_color.rgb = back_ground_color

		#2)Tạo đoạn text bên trong
		tf = txBox.text_frame
		tf.word_wrap = True

		#2.1) Dòng text đầu tiên
		p = tf.paragraphs[0]
		p.text = text_1
		p.font.size = Pt(font_size)
		p.font.color.rgb = RGBColor
		p.font.bold = font_bold
		p.font.name = font_name
		#=== Căn giữa ===
		if alignment == True:
			p.alignment = PP_ALIGN.CENTER

	def create_pie_chart(self,categories,dict_data,left,top,width,height,legend_font=12,chart_title='Title'):
		#1)Thêm categori vào chart
		chart_data = CategoryChartData()
		chart_data.categories = categories
		#1.1)Thêm dữ liệu vào
		for key,value in dict_data.items():
			chart_data.add_series(key, value)

		#2)Thêm chart vào slide
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)
		chart = self.__slide.shapes.add_chart(
				XL_CHART_TYPE.PIE, left, top, width, height, chart_data
		).chart

		#2.1)CHART TITLE
		chart.has_title = True
		chart.chart_title.text_frame.text = chart_title

		#2.2)CHART LEGEND
		chart.has_legend = True
		chart.legend.position = XL_LEGEND_POSITION.BOTTOM
		chart.legend.include_in_layout = False
		chart.legend.font.size = Pt(legend_font)

		#2.3)CÁC SỐ PHẦN TRĂM TRONG CHART
		chart.plots[0].has_data_labels = True
		data_labels = chart.plots[0].data_labels
		data_labels.number_format = '0%'
		data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

	def create_column_stacked_chart(self,categories,dict_data,left,top,width,height,legend_font=12,title='Title',has_legend=True):
		#1)CATEGORIES
		chart_data = CategoryChartData()
		chart_data.categories = categories

		#2)DATA IN CHART
		for key,value in dict_data.items():
			chart_data.add_series(key, value)

		#3)Thêm chart vào slide
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)
		chart = self.__slide.shapes.add_chart(
				XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, chart_data
		).chart
		#3.1)CHART TITLE
		chart.has_title = True
		chart.chart_title.text_frame.text = title

		#3.2)CHART LEGEND
		chart.has_legend = has_legend

		#3.3)NẾU has_legend =True thì mới làm
		if has_legend:
			chart.legend.font.size = Pt(legend_font)
			chart.legend.position = XL_LEGEND_POSITION.TOP
			chart.legend.include_in_layout = False